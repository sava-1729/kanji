"""
A script for reading kanji from PPT slides and saving them in a csv.
"""

import os
import csv
from pptx import Presentation
with open("kanji.csv", "a") as f:
    writer = csv.writer(f)
    for path, dirs, files in os.walk("/path/to/kanji/ppts"):
        for filename in files.keys():
            writer.writerow(["------", filename])
            ppt = Presentation(os.path.join(path, filename))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # print("What is %s?" % shape.text)
                        # yomikata = input()
                        yomikata = ""
                        writer.writerow([shape.text, yomikata])